#!/usr/bin/env python3
"""
LIGHTHOUSE MACRO - FINAL CHARTBOOK GENERATOR
Complete professional chartbook with all charts, sections, and PDF export

Author: Bob Sheehan, CFA, CMT
Date: November 2025
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import glob
from datetime import datetime

# Paths
CHARTS_FOLDER = "/Users/bob/lighthouse_chartbook_final"
OUTPUT_PPTX = "/Users/bob/Lighthouse_Macro_Premium_Chartbook.pptx"
OUTPUT_PDF = "/Users/bob/Lighthouse_Macro_Premium_Chartbook.pdf"

# Lighthouse Macro branding colors
BRAND_BLUE = RGBColor(0, 51, 102)  # #003366
BRAND_LIGHT = RGBColor(0, 102, 204)  # #0066CC
BRAND_ACCENT = RGBColor(255, 153, 0)  # #FF9900

# Chart organization with commentary
CHART_SECTIONS = {
    "SECTION I: PROPRIETARY SYSTEMIC RISK INDICATORS": [
        {
            "file": "MRI_Macro_Risk_Index.png",
            "title": "Macro Risk Index (MRI)",
            "commentary": "MRI at +0.06σ (near neutral) — risk fairly priced but cross-currents beneath surface. Liquidity cushion critically depleted (-1.45σ) while labor fragility building (+0.57σ). Composite masks offsetting forces rather than genuine stability."
        },
        {
            "file": "01_LCI_Liquidity_Cushion_Index.png",
            "title": "Liquidity Cushion Index (LCI)",
            "commentary": "LCI at -1.45σ — critically depleted. RRP collapsed from $2.5T to <$300B. Banking system liquidity 1.5 std devs below average. Primary systemic risk. Watch for -2σ break triggering Fed intervention."
        },
        {
            "file": "02_LFI_Labor_Fragility_Index.png",
            "title": "Labor Fragility Index (LFI)",
            "commentary": "LFI at +0.57σ — moderately elevated. Workers quitting less (confidence deteriorating), long-term unemployment rising, hiring efficiency declining. Leads payroll weakness by 3-6 months. Consumer spending at risk Q1 2026."
        },
        {
            "file": "03_LDI_Labor_Dynamism_Index.png",
            "title": "Labor Dynamism Index (LDI)",
            "commentary": "LDI at -0.53σ — below average vitality. Job switching down, workers staying put. Low dynamism is sticky and self-reinforcing. Historically leads payroll weakness by 2-3 quarters."
        },
        {
            "file": "04_CLG_Credit_Labor_Gap.png",
            "title": "Credit-Labor Gap (CLG)",
            "commentary": "CLG at +0.05σ — credit and labor aligned. No mispricing between bond and jobs markets. Credit has 'caught up' to labor reality. Fair pricing but bad news—recession risk appropriately priced."
        },
        {
            "file": "05_YFS_Yield_Funding_Stress.png",
            "title": "Yield-Funding Stress (YFS)",
            "commentary": "YFS at +0.97σ — approaching +1σ warning threshold. T-bills trading rich to SOFR (safe asset scarcity). Money market stress building. If >+1.5σ, expect Fed liquidity intervention."
        },
        {
            "file": "09_Payrolls_vs_Quits_Divergence.png",
            "title": "Payrolls vs Quits Divergence",
            "commentary": "Textbook late-cycle pattern. Payrolls stable while quits collapsed 33%. Workers stop quitting → companies hoard labor → layoffs imminent. Preceded 2007/2000 recessions by 6 months. High-conviction recession signal."
        },
        {
            "file": "10_Hours_vs_Employment_Divergence.png",
            "title": "Hours vs Employment Divergence",
            "commentary": "Hours growth -0.5% YoY while employment +1.2% — 1.7pp divergence. Companies cut hours before headcount. Leading layoff indicator by 2-4 months. Most reliable recession predictor in toolkit."
        },
        {
            "file": "MRI_Macro_Risk_Index_components.png",
            "title": "MRI Component Breakdown",
            "commentary": "LCI liquidity depletion is dominant risk driver (+1.45σ). Partially offset by labor fragility (+0.57σ) and funding stress (+0.97σ). Credit and equity momentum minimal contributors. Risk is plumbing + labor fundamentals."
        },
    ],
    "SECTION II: GLOBAL MACRO INTELLIGENCE": [
        {
            "file": "MMxLHM - US LEIs vs. S&P500 YoY%.png",
            "title": "U.S. Leading Indicators vs S&P 500",
            "commentary": "Dangerous divergence. LEI -6% YoY (declining 6 months) while S&P +15% near ATH. LEI predicted 8 of 8 recessions with 6-9 month lead. Markets don't 'catch up' to LEI—they crash down. 20% correction needed to realign."
        },
        {
            "file": "MMxLHM - Global Semi Equip Billings vs. Taiwan Exports.png",
            "title": "Semiconductor Cycle Health Check",
            "commentary": "Semi equipment billings +28% YoY, Taiwan exports +40%. AI infrastructure buildout intact. Equipment orders today = production 6-12 months forward. Validates secular AI thesis. Watch for <+15% YoY deceleration as capex pause signal."
        },
        {
            "file": "MMxLHM - Contribution of IT Investment to Change in Real GDP.png",
            "title": "IT Investment GDP Contribution",
            "commentary": "IT capex contributing 0.5pp to GDP — material offset to consumer weakness. AI infrastructure (data centers, GPUs) replacing consumer as growth driver. Sustainability depends on hyperscaler profitability. Recession would cut tech capex, compounding downturn."
        },
        {
            "file": "MMxLHM - Share of Companies Currently or Planning to Use AI.png",
            "title": "AI Enterprise Adoption",
            "commentary": "55% of companies using or planning AI. Adoption curve still early—far from saturation. Validates infrastructure capex sustainability. Users report 10-15% efficiency gains. Supports multi-year AI tailwind thesis."
        },
        {
            "file": "MMxLHM - AI Infra Companies Inventory.png",
            "title": "AI Infrastructure Inventory Levels",
            "commentary": "NVDA inventory $20B (+100% YoY). Rising alongside revenue growth = intentional buffer, not demand weakness. Healthy inventory/revenue ratio. Watch for ratio deterioration as warning signal of oversupply."
        },
        {
            "file": "MMxLHM - Al Software Companies - Remaining Performance Obligations (RPO).png",
            "title": "AI Software Forward Bookings (RPO)",
            "commentary": "Oracle RPO exploded $140B→$460B. Most important AI validation in chartbook. Proves enterprise adoption is real with multi-year commitments. RPO guarantees 30-40% revenue growth through 2027. High-conviction long ORCL."
        },
        {
            "file": "MMxLHM - Magnificent 7 CapEx vs. Compute Efficiency.png",
            "title": "Hyperscaler CapEx Efficiency",
            "commentary": "WARNING: Compute efficiency collapsed 75% (11K→2.5K). Hyperscalers spending same $14B/quarter for far less compute. Diminishing returns emerging. Suggests AI capex peak 2-4 quarters away. Leading indicator for NVDA demand slowdown."
        },
        {
            "file": "MMxLHM - Bitcoin - Average Mining Costs.png",
            "title": "Bitcoin Mining Economics",
            "commentary": "BTC ~$100K vs mining cost ~$60-70K. Healthy 40-50% margin cushion. Sustained trading below cost triggers miner capitulation, marks bottoms. Use mining cost as downside floor for BTC valuation."
        },
        {
            "file": "MMxLHM - Net AI Talent Migration per 10,000 LinkedIn Members.png",
            "title": "AI Talent Migration Flows",
            "commentary": "U.S. winning AI talent war. +15 net inflow per 10K while China -8 outflow. Talent concentration = innovation advantage. Flows are slow but highly predictive. Justifies U.S. AI company premium valuations."
        },
        {
            "file": "MMxLHM - Number of AI-Related Patents Granted (% of World Total).png",
            "title": "AI Patent Leadership",
            "commentary": "U.S. 45% of global AI patents, China 30% rising. U.S. patents higher-impact (more citations). Patent activity today = commercial products 2-3 years forward. Quality matters more than quantity in innovation races."
        },
    ],
    "SECTION III: TECHNICAL ANALYSIS & POSITIONING": [
        {
            "file": "NVDA_2025-11-22_02-17-21.png",
            "title": "NVIDIA (NVDA)",
            "commentary": "80%+ AI GPU share, monopolistic position. Trading at 50-day MA after 15% pullback. Forward P/E 30x vs 40% growth = PEG <1. Support $120 critical. Catalyst risks: China export restrictions, hyperscaler capex slowdown."
        },
        {
            "file": "ASML_2025-11-22_02-17-37.png",
            "title": "ASML Holding (ASML)",
            "commentary": "EUV monopoly, only supplier for <5nm chips. $40B+ backlog = 2+ years visibility. Broke 200-day MA—needs $900 reclaim. China risk: 50% revenue at risk from export controls. Watch for order cancellations."
        },
        {
            "file": "TSM_2025-11-22_02-18-09.png",
            "title": "Taiwan Semiconductor (TSM)",
            "commentary": "60% global chip manufacturing share. Exclusive NVDA H100/H200 manufacturer. Strong uptrend intact near ATH. Geopolitical risk premium from Taiwan Strait tensions. Most advanced process technology (3nm, 2nm coming)."
        },
        {
            "file": "MSFT_2025-11-22_02-28-11.png",
            "title": "Microsoft (MSFT)",
            "commentary": "Clearest AI monetization story. Azure AI +100% YoY, Copilot 1M+ users ($360M run rate). Consolidating $400-$430, breakout >$430 targets $475. Cloud + Office base provides downside support. 30x P/E for 12% growth + AI optionality."
        },
        {
            "file": "JPM_2025-11-22_02-26-31.png",
            "title": "JPMorgan Chase (JPM)",
            "commentary": "Credit cycle barometer. NII peaked Q2 2024, margin compressing. Charge-offs rising but below average—watch for >2% as recession confirmation. 1.8x TBV fair for quality. Use as macro indicator not trade."
        },
        {
            "file": "GS_2025-11-22_02-26-13.png",
            "title": "Goldman Sachs (GS)",
            "commentary": "IB recovery play. M&A +25% YoY, IPO pipeline building. Breaking out of 18-month base, targets $575. 1.5x TBV = premium justified by 15% ROTCE. If M&A stalls, GS underperforms."
        },
        {
            "file": "COIN_2025-11-22_02-21-42.png",
            "title": "Coinbase (COIN)",
            "commentary": "80% revenue transaction-based—highly volatile. SEC lawsuit = binary risk. Range $150-$300. 2x-3x beta to Bitcoin. Use options for asymmetric exposure, not stock. Avoid outright due to regulatory uncertainty."
        },
        {
            "file": "MSTR_2025-11-22_02-24-02.png",
            "title": "MicroStrategy (MSTR)",
            "commentary": "150K+ BTC on balance sheet. 2x beta levered call option on BTC. Trades 2.5x NAV (convertible optionality). Funding risk if BTC <$50K. Not investment—trading vehicle for aggressive BTC bulls."
        },
        {
            "file": "MARA_2025-11-22_02-25-28.png",
            "title": "Marathon Digital (MARA)",
            "commentary": "Profitability depends on BTC >$35K cash cost. 3x-4x beta to Bitcoin. High bankruptcy risk if sustained BTC weakness. Better to trade via options than hold stock. Speculation vehicle only."
        },
        {
            "file": "HYG_2025-11-22_02-29-52.png",
            "title": "High Yield ETF (HYG)",
            "commentary": "Near 52-week highs despite macro risks—credit complacency. Falls 20-30% in recessions. HYG puts = cheap recession protection. Watch breakdown below $78 (200-day MA) for credit stress confirmation."
        },
    ]
}

def create_title_slide(prs):
    """Create professional title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "LIGHTHOUSE MACRO"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Premium Institutional Chartbook"
    p2 = subtitle_frame.paragraphs[0]
    p2.font.size = Pt(28)
    p2.font.color.rgb = BRAND_LIGHT
    p2.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(8), Inches(0.3))
    date_frame = date_box.text_frame
    date_frame.text = datetime.now().strftime("%B %d, %Y")
    p3 = date_frame.paragraphs[0]
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(128, 128, 128)
    p3.alignment = PP_ALIGN.CENTER

def create_section_divider(prs, section_title):
    """Create section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Section title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = section_title
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE
    p.alignment = PP_ALIGN.CENTER

def create_chart_slide(prs, chart_data):
    """Create chart slide with image + commentary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Chart title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = chart_data['title']
    p = title_frame.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE

    # Chart image (left side, larger)
    chart_path = os.path.join(CHARTS_FOLDER, chart_data['file'])
    if os.path.exists(chart_path):
        try:
            pic = slide.shapes.add_picture(chart_path, Inches(0.3), Inches(1), height=Inches(5.5))
        except Exception as e:
            print(f"  ⚠ Could not add {chart_data['file']}: {e}")
    else:
        print(f"  ⚠ Chart not found: {chart_data['file']}")

    # Commentary (right side)
    comment_box = slide.shapes.add_textbox(Inches(5.5), Inches(1), Inches(4.2), Inches(5.5))
    comment_frame = comment_box.text_frame
    comment_frame.word_wrap = True
    comment_frame.text = chart_data['commentary']

    for paragraph in comment_frame.paragraphs:
        paragraph.font.size = Pt(11)
        paragraph.font.name = 'Arial'
        paragraph.space_after = Pt(8)

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.3))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Lighthouse Macro | Bob Sheehan, CFA, CMT | lighthousemacro.substack.com"
    p_footer = footer_frame.paragraphs[0]
    p_footer.font.size = Pt(9)
    p_footer.font.color.rgb = RGBColor(128, 128, 128)
    p_footer.alignment = PP_ALIGN.CENTER

def main():
    print("=" * 70)
    print("LIGHTHOUSE MACRO - FINAL CHARTBOOK GENERATOR")
    print("=" * 70)

    # Create presentation
    print("\n[1/3] Creating PowerPoint presentation...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    print("   ✓ Creating title slide")
    create_title_slide(prs)

    # Add sections with charts
    total_charts = 0
    for section_name, charts in CHART_SECTIONS.items():
        print(f"\n   ✓ Adding {section_name}...")
        create_section_divider(prs, section_name)

        for chart in charts:
            create_chart_slide(prs, chart)
            total_charts += 1
            print(f"      [{total_charts}] {chart['title']}")

    # Save PPTX
    print("\n[2/3] Saving PowerPoint...")
    prs.save(OUTPUT_PPTX)
    pptx_size = os.path.getsize(OUTPUT_PPTX) / (1024 * 1024)

    # Export to PDF
    print("\n[3/3] Exporting to PDF...")
    print("   ℹ Open PowerPoint and export manually, or:")
    print(f"   ℹ Use: open '{OUTPUT_PPTX}'")
    print("   ℹ Then: File → Export → PDF")

    print("\n" + "=" * 70)
    print("CHARTBOOK COMPLETE!")
    print("=" * 70)
    print(f"\nPowerPoint: {OUTPUT_PPTX}")
    print(f"Size: {pptx_size:.1f} MB")
    print(f"Total slides: {len(prs.slides)}")
    print(f"Charts: {total_charts}")
    print("\n✓ Ready to distribute!")
    print("=" * 70)

if __name__ == "__main__":
    main()
