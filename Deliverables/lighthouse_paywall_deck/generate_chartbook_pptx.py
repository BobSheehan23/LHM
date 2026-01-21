#!/usr/bin/env python3
"""
LIGHTHOUSE MACRO - AUTOMATED PPTX CHARTBOOK GENERATOR
Creates professional PowerPoint with one chart + commentary per slide

Author: Bob Sheehan, CFA, CMT
Date: November 2025
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import re
from datetime import datetime

# Paths
CHARTS_FOLDER = "/Users/bob/lighthouse_chartbook_final"
COMMENTARY_FILE = f"{CHARTS_FOLDER}/CHARTBOOK_GUIDE.md"
OUTPUT_PPTX = "/Users/bob/Lighthouse_Macro_Premium_Chartbook.pptx"

# Lighthouse Macro branding colors
BRAND_BLUE = RGBColor(0, 51, 102)  # #003366
BRAND_LIGHT = RGBColor(0, 102, 204)  # #0066CC
BRAND_ACCENT = RGBColor(255, 153, 0)  # #FF9900

def parse_commentary():
    """Parse markdown commentary into structured data"""
    with open(COMMENTARY_FILE, 'r') as f:
        content = f.read()

    # Extract chart sections with ### headers
    chart_sections = re.findall(
        r'### (.*?\.png)\n\n\*\*(.*?)\*\*\n\n(.*?)(?=\n\n---|$)',
        content,
        re.DOTALL
    )

    charts = []
    for filename, title, commentary in chart_sections:
        # Clean up filename
        filename = filename.strip()
        # Clean up commentary (remove extra markdown)
        commentary = commentary.replace('**', '')

        chart_path = os.path.join(CHARTS_FOLDER, filename)
        if os.path.exists(chart_path):
            charts.append({
                'filename': filename,
                'title': title.strip(),
                'commentary': commentary.strip(),
                'path': chart_path
            })

    return charts

def create_title_slide(prs):
    """Create cover slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "LIGHTHOUSE MACRO"

    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(left, top + Inches(1), width, Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Premium Institutional Chartbook"

    p2 = subtitle_frame.paragraphs[0]
    p2.font.size = Pt(28)
    p2.font.color.rgb = BRAND_LIGHT
    p2.alignment = PP_ALIGN.CENTER

    # Add date
    date_box = slide.shapes.add_textbox(left, top + Inches(1.8), width, Inches(0.3))
    date_frame = date_box.text_frame
    date_frame.text = datetime.now().strftime("%B %d, %Y")

    p3 = date_frame.paragraphs[0]
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(128, 128, 128)
    p3.alignment = PP_ALIGN.CENTER

    return slide

def create_chart_slide(prs, chart_data):
    """Create slide with chart image + commentary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add chart title at top
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = chart_data['title']
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE

    # Add chart image (left side, larger)
    try:
        left = Inches(0.3)
        top = Inches(1)
        height = Inches(5.5)
        pic = slide.shapes.add_picture(chart_data['path'], left, top, height=height)
    except Exception as e:
        print(f"  ⚠ Could not add image {chart_data['filename']}: {e}")

    # Add commentary (right side)
    comment_left = Inches(5.5)
    comment_top = Inches(1)
    comment_width = Inches(4.2)
    comment_height = Inches(5.5)

    comment_box = slide.shapes.add_textbox(comment_left, comment_top, comment_width, comment_height)
    comment_frame = comment_box.text_frame
    comment_frame.word_wrap = True

    # Split commentary into paragraphs for better formatting
    paragraphs = chart_data['commentary'].split('\n\n')

    for i, para_text in enumerate(paragraphs):
        if i > 0:
            comment_frame.add_paragraph()
        p = comment_frame.paragraphs[i]
        p.text = para_text.strip()
        p.font.size = Pt(10)
        p.font.name = 'Arial'
        p.space_after = Pt(6)

        # Highlight key phrases
        if any(keyword in para_text for keyword in ['Action:', 'Trade:', 'Watch:', 'Critical:']):
            p.font.bold = True
            p.font.color.rgb = BRAND_ACCENT

    # Add footer with source
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.3))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Lighthouse Macro | Bob Sheehan, CFA, CMT | lighthousemacro.substack.com"

    p_footer = footer_frame.paragraphs[0]
    p_footer.font.size = Pt(9)
    p_footer.font.color.rgb = RGBColor(128, 128, 128)
    p_footer.alignment = PP_ALIGN.CENTER

    return slide

def main():
    print("=" * 70)
    print("LIGHTHOUSE MACRO - PPTX CHARTBOOK GENERATOR")
    print("=" * 70)

    # Parse commentary
    print("\n[1/3] Parsing commentary markdown...")
    charts = parse_commentary()
    print(f"   ✓ Found {len(charts)} charts with commentary")

    # Create presentation
    print("\n[2/3] Creating PowerPoint presentation...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add title slide
    print("   ✓ Creating title slide")
    create_title_slide(prs)

    # Add chart slides
    print("   ✓ Adding chart slides...")
    for i, chart in enumerate(charts, 1):
        create_chart_slide(prs, chart)
        print(f"      [{i}/{len(charts)}] {chart['filename']}")

    # Save
    print("\n[3/3] Saving presentation...")
    prs.save(OUTPUT_PPTX)

    file_size = os.path.getsize(OUTPUT_PPTX) / (1024 * 1024)

    print("\n" + "=" * 70)
    print("CHARTBOOK COMPLETE!")
    print("=" * 70)
    print(f"\nOutput: {OUTPUT_PPTX}")
    print(f"Size: {file_size:.1f} MB")
    print(f"Slides: {len(charts) + 1} (1 title + {len(charts)} charts)")
    print("\n✓ Ready to distribute to premium subscribers!")
    print("=" * 70)

if __name__ == "__main__":
    main()
